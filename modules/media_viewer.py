"""
LMS Media Viewer & Document Parser Module.
Provides native parsing, inspection, and playback for PDFs, PPTX slide decks,
and streaming video recordings within the Streamlit workspace.
"""

import os
from typing import Optional
import pypdf
from pptx import Presentation
import streamlit as st


def render_pdf(file_path_or_buffer: Any) -> None:
    """Extracts text and renders a page-by-page navigator for PDF lecture notes and handbooks."""
    try:
        reader = pypdf.PdfReader(file_path_or_buffer)
        total_pages = len(reader.pages)

        if total_pages == 0:
            st.warning("The selected PDF document contains no readable pages.")
            return

        st.info(f"📄 Document Loaded: **{total_pages} Pages**")
        page_num = st.number_input("Select Page", min_value=1, max_value=total_pages, value=1, step=1)
        
        selected_page = reader.pages[page_num - 1]
        text_content = selected_page.extract_text() or ""

        if text_content.strip():
            st.text_area(f"Extracted Page Content (Page {page_num})", text_content, height=350)
        else:
            st.info(f"Page {page_num} contains diagrams or scanned images without embedded text.")
    except Exception as e:
        st.error(f"Error reading PDF file: {str(e)}")


def render_ppt(file_path_or_buffer: Any) -> None:
    """Parses and renders slide text elements from PowerPoint presentations (.pptx)."""
    try:
        prs = Presentation(file_path_or_buffer)
        total_slides = len(prs.slides)

        if total_slides == 0:
            st.warning("The selected presentation contains no slides.")
            return

        st.info(f"📊 Presentation Loaded: **{total_slides} Slides**")
        slide_num = st.number_input("Select Slide", min_value=1, max_value=total_slides, value=1, step=1)
        
        target_slide = prs.slides[slide_num - 1]
        extracted_texts = []

        for shape in target_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        extracted_texts.append(line)

        st.markdown(f"**Slide {slide_num} Structure & Key Points:**")
        if extracted_texts:
            for item in extracted_texts:
                st.markdown(f"- {item}")
        else:
            st.info(f"Slide {slide_num} contains graphics or charts without direct text frames.")
    except Exception as e:
        st.error(f"Error parsing PowerPoint presentation: {str(e)}")


def render_video(file_path_or_url: str) -> None:
    """Streams live session recordings or video lecture archives."""
    if not file_path_or_url:
        st.warning("Please provide a valid video link or file path.")
        return

    try:
        st.video(file_path_or_url)
    except Exception as e:
        st.error(f"Error loading video stream: {str(e)}")
